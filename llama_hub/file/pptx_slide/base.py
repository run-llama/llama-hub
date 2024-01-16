"""Read Microsoft PowerPoint files."""

from pathlib import Path
from typing import Dict, List, Optional

from llama_index.readers.base import BaseReader
from llama_index.readers.schema.base import Document


class PptxSlideReader(BaseReader):
    """Powerpoint Slides Reader.

    Create a list of Documents corresponding to the Slides of the presentation.

    """

    def __init__(self) -> None:
        """Init reader."""

    def load_data(
        self,
        file: Path or str,
        extra_info: Optional[Dict] = None,
    ) -> List[Document]:
        """Load pptx file to create slide Documents"""
        from pptx import Presentation

        if isinstance(file, str):
            file_path = Path(file)
        else:
            file_path = file

        presentation = Presentation(file_path)
        slide_docs = [
            Document(
                text="\n".join(
                    [shape.text for shape in slide.shapes if hasattr(shape, "text")]
                ),
                extra_info={
                    "source": file_path.name,
                    "shapes": [
                        {
                            "text": shape.text,
                            "name": shape.name,
                            "shape_id": shape.shape_id,
                            "shape_type": shape.shape_type,
                        }
                        for shape in slide.shapes
                        if hasattr(shape, "text")
                    ],
                },
            )
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        ]

        return slide_docs
